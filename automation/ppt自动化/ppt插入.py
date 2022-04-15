from pptx import Presentation
from pptx.util import Inches,Pt
ppt=Presentation()
#在ppt中插入一个幻灯片
slide=ppt.slides.add_slide(ppt.slide_layouts[6])
#插入图片
left=Inches(2)
top=Inches(2)
width=Inches(3)
height=Inches(3)
textbox=slide.shapes.add_picture('E:\图片\壁纸\女孩喝啤酒 飘窗 城市夜景 雨天 è 好看唯美4k动漫壁纸_彼岸图网.jpg',left,top,width,height)

#插入表格
rows=2
cols=2
left=Inches(1)
top=Inches(1)
width=Inches(4)
height=Inches(4)
table=slide.shapes.add_table(rows,cols,left,top,width,height).table
table.columns[0].width=Inches(1)
table.columns[1].width=Inches(3)
table.cell(0,0).text='55'
table.cell(0,1).text='425'
table.cell(1,0).text='54165'
table.cell(1,1).text='6546'

ppt.save('insert.pptx')