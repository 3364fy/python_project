from pptx import Presentation
from pptx.util import Inches,Pt
ppt=Presentation()
#在ppt中插入一个幻灯片
slide=ppt.slides.add_slide(ppt.slide_layouts[1])
#文本框
body_shape=slide.shapes.placeholders
body_shape[0].text='这是占位符【0】'
body_shape[1].text='这是占位符【1】'

# title_shape=slide.shapes.title
# title_shape.text='这里是标题'
# subtitle=slide.shapes.placeholders[1]
# subtitle.text='这里是文本框'

#在文本框内写入段落
new_paragraph=body_shape[1].text_frame.add_paragraph()
new_paragraph.text='新段落'
new_paragraph.font.bold=True
new_paragraph.font.italic=True
new_paragraph.font.size=Pt(15)
new_paragraph.font.underline=True

#新建文本框
left=Inches(2)
top=Inches(2)
width=Inches(3)
height=Inches(3)
textbox=slide.shapes.add_textbox(left,top,width,height)
textbox.text='这是新文本框'
new_para=textbox.text_frame.add_paragraph()
new_para.text='这是新文本框的第二段'

ppt.save('test.pptx')